# converter/views.py

from django.shortcuts import render, redirect
from django.http import HttpResponse
from .forms import ImageUploadForm
from .models import Conversion
from PIL import Image
from io import BytesIO
import os
from django.conf import settings
from django.core.files.storage import FileSystemStorage
from docx import Document
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation




def home(request):
    if request.method == 'POST':
        form = ImageUploadForm(request.POST, request.FILES)
        if form.is_valid():
            conversion = form.save(commit=False)

            # Convert image to PDF
            image_file = request.FILES['image']
            image = Image.open(image_file)

            # Create a BytesIO object to hold the PDF data
            pdf_bytes = BytesIO()
            pdf_image = image.convert('RGB')
            pdf_image.save(pdf_bytes, format='PDF')
            pdf_bytes.seek(0)  # Rewind the BytesIO object to the beginning

            # Set the PDF file name
            pdf_filename = f"{os.path.splitext(image_file.name)[0]}.pdf"

            # Return the PDF as an HttpResponse
            response = HttpResponse(pdf_bytes.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{pdf_filename}"'
            return response

    else:
        form = ImageUploadForm()
    return render(request, 'index.html', {'form': form})



def convert_to_pdf(request):
    if request.method == 'POST':
        form = ImageUploadForm(request.POST, request.FILES)
        if form.is_valid():
            conversion = form.save(commit=False)

            # Convert image to PDF
            image_file = request.FILES['image']
            image = Image.open(image_file)

            # Create a BytesIO object to hold the PDF data
            pdf_bytes = BytesIO()
            pdf_image = image.convert('RGB')
            pdf_image.save(pdf_bytes, format='PDF')
            pdf_bytes.seek(0)  # Rewind the BytesIO object to the beginning

            # Set the PDF file name
            pdf_filename = f"{os.path.splitext(image_file.name)[0]}.pdf"

            # Return the PDF as an HttpResponse
            response = HttpResponse(pdf_bytes.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{pdf_filename}"'
            return response

    else:
        form = ImageUploadForm()
    return render(request, 'index.html', {'form': form})







def doc_to_pdf_view(request):
    if request.method == 'POST' and request.FILES['doc_file']:
        uploaded_file = request.FILES['doc_file']

        try:
            # Create an in-memory file for the PDF
            pdf_buffer = BytesIO()
            
            # Convert DOCX to PDF
            convert_docx_to_pdf(uploaded_file, pdf_buffer)

            # Return the PDF as a response
            pdf_buffer.seek(0)
            response = HttpResponse(pdf_buffer, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="converted_file.pdf"'
            return response

        except Exception as e:
            return HttpResponse(f"Conversion failed: {str(e)}", status=500)

    return render(request, 'dic_to_pdf.html')

def convert_docx_to_pdf(input_file, pdf_buffer):
    doc = Document(input_file)

    pdf = canvas.Canvas(pdf_buffer, pagesize=letter)

    width, height = letter
    x, y = 72, height - 72  

    for paragraph in doc.paragraphs:
        text = paragraph.text
        if y <= 72:  # Create a new page if we reach the bottom of the page
            pdf.showPage()
            y = height - 72 

        pdf.drawString(x, y, text)
        y -= 15  # Adjust line height

    pdf.save()






def ppt_to_pdf_view(request):
    if request.method == 'POST' and request.FILES['ppt_file']:
        uploaded_file = request.FILES['ppt_file']

        try:
            pdf_buffer = BytesIO()
            convert_pptx_to_pdf(uploaded_file, pdf_buffer)
            
            pdf_buffer.seek(0)
            response = HttpResponse(pdf_buffer, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="converted_ppt.pdf"'
            return response

        except Exception as e:
            return HttpResponse(f"Conversion failed: {str(e)}", status=500)

    return render(request, 'ppt_to_pdf.html')

def convert_pptx_to_pdf(input_file, pdf_buffer):
    prs = Presentation(input_file)
    pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
    width, height = letter

    for slide in prs.slides:
        y = height - 50
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                pdf.drawString(72, y, text)
                y -= 20

        pdf.showPage()

    pdf.save()




import io
from django.http import HttpResponse
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import openpyxl

def excel_to_pdf_view(request):
    if request.method == 'POST' and request.FILES['excel_file']:
        uploaded_file = request.FILES['excel_file']

        try:
            # Create a BytesIO object to hold the PDF in memory
            pdf_buffer = io.BytesIO()

            # Convert Excel to PDF (write to pdf_buffer instead of saving to a file)
            convert_excel_to_pdf(uploaded_file, pdf_buffer)

            # Get the PDF content and prepare the response
            pdf_buffer.seek(0)
            response = HttpResponse(pdf_buffer, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="converted_excel.pdf"'
            return response

        except Exception as e:
            return HttpResponse(f"Conversion failed: {str(e)}", status=500)

    return render(request, 'excel_to_pdf.html')

def convert_excel_to_pdf(input_file, pdf_buffer):
    wb = openpyxl.load_workbook(input_file)
    ws = wb.active

    # Create a new PDF file in the BytesIO buffer
    pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
    width, height = letter
    x, y = 50, height - 50

    # Loop through rows and columns and write the data to the PDF
    for row in ws.iter_rows(values_only=True):
        row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
        
        if y <= 50:  # Create a new page if we reach the bottom
            pdf.showPage()
            y = height - 50
        
        pdf.drawString(x, y, row_text)
        y -= 15  # Move to the next line

    pdf.save()  # Finalize the PDF



import os
from django.shortcuts import render
from django.http import HttpResponse
from xhtml2pdf import pisa

def html_to_pdf_view(request):
    if request.method == 'POST' and request.FILES['html_file']:
        uploaded_file = request.FILES['html_file']
        html_content = uploaded_file.read().decode('utf-8')  # Read and decode the uploaded file

        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="converted.pdf"'
        pisa_status = pisa.CreatePDF(html_content, dest=response)

        if pisa_status.err:
            return HttpResponse('Error rendering PDF', status=400)
        return response

    return render(request, 'html_to_pdf.html')






from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse
from django.core.files.storage import FileSystemStorage
from PyPDF2 import PdfFileReader, PdfFileWriter
import io



# def upload_pdf(request):
#     if request.method == 'POST' and request.FILES['pdf_file']:
#         pdf_file = request.FILES['pdf_file']
#         fs = FileSystemStorage()
#         filename = fs.save(pdf_file.name, pdf_file)
#         uploaded_file_url = fs.url(filename)
#         return render(request, 'edit_pdf.html', {'pdf_url': uploaded_file_url, 'pdf_name': pdf_file.name})
#     return render(request, 'upload_pdf.html')

# def edit_pdf(request):
#     if request.method == 'POST' and request.FILES['pdf_file']:
#         pdf_file = request.FILES['pdf_file']
#         fs = FileSystemStorage()
#         filename = fs.save(pdf_file.name, pdf_file)

#         # Here, we're assuming you're editing the PDF. For now, this just sends the file back unchanged.
#         response = HttpResponse(content_type='application/pdf')
#         response['Content-Disposition'] = f'attachment; filename="edited_{pdf_file.name}"'
        
#         # Send back the original file without modifications (you can change this to edit the file)
#         with open(fs.path(filename), 'rb') as pdf:
#             response.write(pdf.read())
#         return response
#     return render(request, 'edit_pdf.html')

def download_pdf(request, pk):
    pdf_file = get_object_or_404(PdfFile, pk=pk)
    response = HttpResponse(pdf_file.pdf_file, content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="{pdf_file.name}"'
    return response


from django.shortcuts import render
from django.core.files.storage import FileSystemStorage
from django.http import HttpResponse
import os

from django.shortcuts import render
from django.core.files.storage import FileSystemStorage
from django.http import HttpResponse


# def edit_pdf(request):
#     print('\n ➡ 296 request:', request)
#     if request.method == 'POST' and request.FILES.get('pdf_file'):
#         pdf_file = request.FILES['pdf_file']
#         fs = FileSystemStorage()
#         filename = fs.save(pdf_file.name, pdf_file)

#         # URL where the PDF is stored
#         pdf_url = fs.url(filename)
#         print('\n ➡ 303 pdf_url:', pdf_url)
    
#         pdf_url = "http://127.0.0.1:8000" + pdf_url
#         # print('\n ➡ 306 pdf_url:', pdf_url)
        
#         return render(request, 'edit_pdf.html', {'pdf_url': pdf_url, 'pdf_name': filename})

#         # For now, send the original file back unchanged
#         response = HttpResponse(content_type='application/pdf')
#         response['Content-Disposition'] = f'attachment; filename="edited_{pdf_file.name}"'
        

#         with open(fs.path(filename), 'rb') as pdf:
#             response.write(pdf.read())
#         return response
    
    # Initial page load - No file yet, so no PDF URL
    # return render(request, 'edit_pdf.html', {'pdf_url': None, 'pdf_name': None})



from django.core.files.storage import FileSystemStorage
from django.http import HttpResponse
from django.shortcuts import render
from PyPDF2 import PdfFileReader, PdfFileWriter
import os

def edit_pdf(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        # When the user uploads a PDF
        pdf_file = request.FILES['pdf_file']
        fs = FileSystemStorage()
        filename = fs.save(pdf_file.name, pdf_file)

        # URL where the PDF is stored
        pdf_url = fs.url(filename)
        pdf_path = fs.path(filename)

        # Send back the rendered page with the PDF embedded
        return render(request, 'edit_pdf.html', {'pdf_url': pdf_url, 'pdf_name': pdf_file.name})

    elif request.method == 'GET':
        # Initial GET request, no PDF yet, just render the form
        return render(request, 'upload_pdf.html')
